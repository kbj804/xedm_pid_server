'''
pip install pdfminer.six  # https://pdfminersix.readthedocs.io/en/latest/tutorial/extract_pages.html
pip install 'olefile'  or '-U olefile'
pip install python-pptx
pip install python-docx
pip install openpyxl
pip install xmltodict
'''
import os
# from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
# from pdfminer.converter import HTMLConverter
# from pdfminer.converter import TextConverter
# from pdfminer.layout import LAParams
# from pdfminer.pdfpage import PDFPage
# from io import StringIO
# import olefile
from pptx import Presentation
from docx import Document
import csv
import openpyxl
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer
import  json


from common.consts import SAMPLE_FOLDER_PATH, UPLOAD_DIRECTORY, IMG_OUTPUT_PATH
from utils.logger_handler import get_logger
from utils.file_module.pdf_layout_scanner import get_pages

logger = get_logger()

class loadFileManager:
    # 파일 이름, 확장자 분리
    def __init__(self, path):
        # dir_path = os.getcwd()
        self.path = path
        basename = os.path.basename(self.path)
        self.name, self.dotext = os.path.splitext(basename)
        self.ext = self.dotext.replace(".",'',1)

        self.data = None
        """
        [
            {
                page:1
                td: dfsdf
            },{
                page:2
                td: asdasdasda
            }
        ]
        """

        # 읽을 수 있는 파일인지 검사하고 읽을 수 있으면 data입력
        if self.check_ext():
            self.data = self.read_file()

        else:
            logger.info(f'### Error ### Load_file_manager.py 에서 지원하지 않는 확장자 입니다. : {self.name}.{self.ext}')
            print(f"### load_file_manager 에서 지원하지 않는 확장자 입니다. ### {self.name}.{self.ext}")
            pass
            
    # 확장자에 맞는 read 함수로 매핑
    def read_file(self):
        result = self.read_function[self.ext](self)
        return result

    # 읽을 수 있는 확장자인가 검사
    def check_ext(self):
        try:
            if self.read_function[self.ext]:
                return True
            else:
                print("@ # @ # {} For Debugging... @ # @ #".format(self.ext))
                return False

        except Exception as e:
            print("####ERROR#### {0} dose not exist in ext Dictionary".format(e))

    # 확장자 별 데이터 오픈
    # pdf, hwp, ppt, docx, ....

    def read_pdf(self):
        # using pdfminer.six  / No pdfminer
        result = []
        for i, page_layout in enumerate(extract_pages(self.path)):
            page_contents=''
            for element in page_layout:
                if isinstance(element, LTTextContainer):
                    page_contents += element.get_text()
            result.append({"page":i, "td":page_contents})
        self.read_pdf_for_image_extract()
        
        return result
        

        # pdf 추출하는 다른 방법 (pdfminer)
        '''rsrcmgr = PDFResourceManager()
        retstr = StringIO()
        codec = 'utf-8'
        laparams = LAParams()
        # f = open('./out.html', 'wb')
        # device = HTMLConverter(rsrcmgr, f, codec=codec, laparams=laparams)
        device = TextConverter(rsrcmgr, retstr, codec=codec, laparams=laparams)
        
        with open(self.path, 'rb') as fp:
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            # password = ""
            maxpages = 0 #is for all
            # caching = True
            pagenos=set()
            
            for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages, check_extractable=True):
                interpreter.process_page(page)
            
            # 이건 전체 페이지 전부다 리턴할때 쓰는거
            str = retstr.getvalue()

            fp.close()
        
        device.close()
        retstr.close()
        f.close()
    
        return str
        '''

    def read_hwp(self):
        """
        hwp -> pdf 로 변환 후 pdf_read를 해서 값을 반환해 주기 때문에 약간 골치아픔..
        window에서 RegisterMoudle 설정해줘야 하고, 리눅스에서는 어떻게 될 지 모르겠음.
        """
        import win32com.client as win32
        import win32gui
        import win32con

        hwp = win32.gencache.EnsureDispatch("HWPFrame.HwpObject")  # 한/글 열기
        # hwp.XHwpWindows.Item(0).Visible = True  # 숨김해제
        try:
            if hwp.RegisterModule("FilePathCheckDLL", "FilePathChekcerModule"):  # 보안모듈 실행. (경고창 없어짐)
                hwp.Clear(option=1)
                # hwnd = win32gui.FindWindow(None, "빈 문서 1 - 한글")
                hwp.Open(self.path)  # 한/글파일 열기
                # win32gui.ShowWindow(hwnd, win32con.SW_HIDE)

                # current_page = hwp.XHwpDocuments.Item(0).XHwpDocumentInfo.CurrentPage
                # print(current_page)
                pdf_path = SAMPLE_FOLDER_PATH+self.name+".pdf"
                self.path = pdf_path
                hwp.SaveAs(pdf_path, "PDF")

                return self.read_pdf()
                # hwp.Quit()

        except Exception as e:
            print("#### ERROR #### ",e)

        finally:
            print("## QUIT HWP ###")
            if os.path.isfile(self.path):
                os.remove(self.path)
            hwp.Quit()


    def read_hwp_text(self):
        import olefile 
        result :list = []
        f = olefile.OleFileIO(self.path)
        encoded_text = f.openstream('PrvText').read()
        decoded_text = encoded_text.decode('UTF-16')

        result.append({"page":0, "td": decoded_text })
        return result
        


    def read_pptx(self):
        """
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        Doc : https://python-pptx.readthedocs.io/en/latest/api/slides.html#slides-objects
        """
        pptxdoc = Presentation(self.path)
        result :list = []
        

        for pn, slide in enumerate(pptxdoc.slides):
            text_runs = []
            table_runs = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    # 이게 왜 필요하지
                    # if not shape.has_text_frame:
                    #     continue
                    for paragraph in shape.text_frame.paragraphs:
                        # 1. run -> 모든 내용 제목 까지 전부다
                        for run in paragraph.runs:
                            text_runs.append(run.text)

                        # 2. paragraph -> 페이지별로 구분됨
                        # text_runs.append(paragraph.text)

                elif shape.has_table:
                    # print("table")
                    tb1 = shape.table
                    row_count = len(tb1.rows)
                    col_count = len(tb1.columns)
                    for r in range(0, row_count):
                        for c in range(0, col_count):
                            cell = tb1.cell(r,c)
                            paragraphs = cell.text_frame.paragraphs 
                            for paragraph in paragraphs:
                                for run in paragraph.runs:
                                    table_runs.append(run.text)
                        table_runs.append('\n')
                else:
                    # 기타 이미지나 뭐 그런거,,
                    pass
            # print(text_runs))
            result.append({"page":pn, "td": ' '.join(text_runs), "table": '|'.join(table_runs) })

        # result.append({"page":0, "td": '\n'.join(answer) })
        # print(' '.join(text_runs))
        return result

    def read_docx(self):
        """
        Docx File 변환인데, TEXT, TABLE은 어떻게든 가능하나 Page를 뽑는게 안됨
        """
        result=[]
        document = Document(self.path)
        fullText = []
        pn = 1
        for para in document.paragraphs:
            for run in para.runs:
                fullText.append(run.text)
                if 'lastRenderedPageBreak' in run._element.xml:
                    pn += 1
                    print(f"{pn}page -> text: {run.text}")
                    # print(run._element.xml)
        result.append({"page":0, "td": '\n'.join(fullText) })
        return result

    # def read_docx(self):
        """Linux에서 지원 안됨 docx2pdf MSoffice를 설치해야 됨"""
    #     from docx2pdf import convert
    #     output = UPLOAD_DIRECTORY + "output.pdf"
    #     convert(self.path, output)
    #     self.path = output
    #     return self.read_pdf()


    # csv 코덱문제 해결하고 엑셀읽기 하고 html읽기 하면됨 
    def read_csv(self):
        result=[]
        read_list=[]
        with open(self.path, 'r', encoding='UTF8') as file:
            file_read = csv.reader(file)
            for line in file_read:
                read_list.append(line)

            # 2차원 리스트 -> 1차원으로
            answer = sum(read_list,[])

        result.append({"page":0, "td": '\n'.join(answer) })

        return result

    # 리스트 값 안에 숫자형도 있어서 \n . join 사용이 안됨
    # 리스트가 3차원 까지 가서 코드가 다소 복잡함
    # 마지막 리스트가 문자열로 치환이 안되어 추수 수정 필요
    def read_xlsx(self):
        workbook = openpyxl.load_workbook(self.path, data_only=True)
        # Sheet 목록
        sheet_list = workbook.sheetnames
        # all_sheet_value=[]
        
        result =[]
        
        # Sheet 별 탐색
        for i, sheet in enumerate(sheet_list):
            all_values = []
            workSheet = workbook[sheet]
            for row in workSheet.rows:
                row_value =[]
                for cell in row:
                    if cell.value is None:
                        pass
                    else:
                        row_value.append(cell.value)
                all_values.append(row_value)
            # 차원축소1
            answer = sum(all_values,[])
            
            # answer가 list인데, 내부 요소중에 int type이 있어서 그냥 join이 불가능함. 그래서 list(map) 함수로 int를 문자로 바꾼 후 작업
            result.append({"page":i, "td": ' '.join(list(map(str,answer)))})

        #     all_sheet_value.append(answer)
        # # 차원축소2
        # answer2 = sum(all_sheet_value, [])

        # 리스트 강제로 str로 함 ㅠ
        return result
        

    def read_txt(self):
        result : list=[]
        with open(self.path, 'r', encoding='UTF8') as f:
            data = f.read()
            result.append({"page":0, "td": data })
        return result
    
    def read_html(self):
        pass

    def read_json(self):
        """ json parser
        """
        result : list=[]
        with open(self.path, 'r') as json_file:
            json_data = json.load(json_file) # type: dict
            # type(json.dumps(json_data)) # type: str
            result.append({"page":0,  "td": json_data})
        
        return result
    
    def read_xml(self):
        """ xml parser
        : return dict
        """
        from xml.etree.ElementTree import parse, fromstring
        import xmltodict
        result : list = []

        with open(self.path) as fd:
            doc = xmltodict.parse(fd.read(), process_namespaces=True)
        
        data = json.dumps(doc)
        # tree = parse(self.path)
        # root = tree.getroot()
        # print(tree)
        # for child in root.iter():
        #     print(child.tag, child.attrib)

        result.append({"page":0,  "td": data})

        return result

    def read_pdf_for_image_extract(self):
        """OCR 기능 연동을 위한 이미지 추출 모듈"""
        # pip install pymuPDF pillow
        import fitz
        import io
        import os
        from PIL import Image

        with fitz.open(self.path) as pdf_file :

            # iterate over PDF pages
            for page_index in range(len(pdf_file)):
                # get the page itself
                page = pdf_file[page_index]
                image_list = page.getImageList()
                # printing number of images found in this page
                if image_list:
                    print(f"[+] Found a total of {len(image_list)} images in page {page_index}")
                else:
                    print("[!] No images found on page", page_index)
                for image_index, img in enumerate(page.getImageList(), start=1):
                    # get the XREF of the image
                    xref = img[0]
                    # extract the image bytes
                    base_image = pdf_file.extractImage(xref)
                    image_bytes = base_image["image"]
                    # get the image extension
                    image_ext = base_image["ext"]
                    # load it to PIL
                    image = Image.open(io.BytesIO(image_bytes))
                    # save it to local disk
                    IMG_PATH = os.path.join(IMG_OUTPUT_PATH, f"image{page_index+1}_{image_index}.{image_ext}")
                    image.save(open(IMG_PATH, "wb"))
        self.clear_img_folder()
            
    def clear_img_folder(self):
        """clear img extraction output folder"""
        file_list = os.listdir(IMG_OUTPUT_PATH) 
        for file in file_list:
            RM_PATH = os.path.join(IMG_OUTPUT_PATH, file)
            os.remove(RM_PATH)

    read_function = {
        'pdf': read_pdf,
        'hwp': read_hwp_text,
        'pptx': read_pptx,
        'docx': read_docx,
        'csv': read_csv,
        'xlsx': read_xlsx,
        'txt': read_txt,
        # 'html': read_html,
        'json': read_json,
        'xml': read_xml
    }


# a = loadFileManager(SAMPLE_FOLDER_PATH + 'docx_sample4.docx')
# a = loadFileManager(SAMPLE_FOLDER_PATH + 'xml_sample.xml')
# a = loadFileManager(SAMPLE_FOLDER_PATH + 'json_sample.json')

# print(a.data)
